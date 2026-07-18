"""Generate the TORQ pitch deck (.pptx).

Run:  uv run python scripts/build_deck.py
Output: data-room/TORQ_deck.pptx
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.stdout.reconfigure(encoding="utf-8")

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "data-room" / "TORQ_deck.pptx"

BG = RGBColor(0x0F, 0x17, 0x20)
CARD = RGBColor(0x11, 0x1C, 0x28)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTE = RGBColor(0x8B, 0x98, 0xA5)

W, H = Inches(13.333), Inches(7.5)


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _box(slide, x, y, w, h, size, text, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return tb


def _accent(slide):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.18), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()


def content(prs, title, bullets, kicker=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent(s)
    if kicker:
        _box(s, Inches(0.7), Inches(0.5), Inches(11), Inches(0.4), 14, kicker.upper(), GREEN, True)
    _box(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0), 34, title, TEXT, True)
    body = "\n".join(f"•  {b}" for b in bullets)
    _box(s, Inches(0.8), Inches(2.1), Inches(11.8), Inches(5), 20, body, TEXT)
    return s


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1. Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, Inches(0.9), Inches(2.6), Inches(11), Inches(1.2), 60, "TORQ", GREEN, True)
    _box(s, Inches(0.95), Inches(3.7), Inches(11), Inches(0.8), 28, "Fault-to-Fix Engine", TEXT, True)
    _box(s, Inches(0.95), Inches(4.5), Inches(11.5), Inches(0.8), 20,
         "From fault code to fixed: autonomous machine diagnosis and repair dispatch.", MUTE)
    _box(s, Inches(0.95), Inches(6.4), Inches(11), Inches(0.6), 16,
         "Automate or Die Hackathon  ·  Team La Scaloneta", MUTE)

    content(prs, "The most expensive 30 minutes in manufacturing", [
        "Unplanned downtime costs hundreds of dollars per minute.",
        "When a machine faults, diagnosis is still manual: ~30 min digging through 500-page manuals while the line is stopped.",
        "Expert know-how walks out the door; the same fault is re-diagnosed from scratch every time.",
        "Dispatch is improvised: 15-60 min of phone calls to find the right technician.",
    ], kicker="Problem")

    content(prs, "A zero-touch fault-to-fix pipeline, not a chatbot", [
        "Machine faults (MQTT) trigger an AI agent automatically.",
        "It reads plant manuals + repair history via MCP, on-premise.",
        "Generates a trilingual (FR/AR/EN) work order: steps, parts, safety.",
        "One-tap supervisor approval, then skill-matched dispatch to the on-shift technician's phone.",
        "Repair outcomes feed back, so every future diagnosis gets better.",
    ], kicker="Solution")

    content(prs, "How it works", [
        "Fault event  ->  Diagnosis agent (DeepSeek + retrieval)",
        "MCP knowledge server: manuals + history, data stays local",
        "Work order + trilingual PDF  ->  supervisor approval queue",
        "Skill-matched routing  ->  WhatsApp / SMS dispatch (Twilio)",
        "Outcome feedback  ->  self-improving knowledge base",
        "Live downtime dashboard: time-to-diagnosis, MTTR, resolution rate",
    ], kicker="Architecture")

    content(prs, "What makes TORQ different", [
        "Prediction to autonomous action: it diagnoses AND dispatches, the human only approves.",
        "MCP as an on-premise trust layer: proprietary manuals never leave the plant.",
        "Human-in-the-loop by design: speed of an agent, accountability of a manager.",
        "Self-improving: tribal knowledge becomes a durable plant asset.",
        "Trilingual (FR/AR/EN) for the real North African shop floor.",
    ], kicker="Innovation")

    # Impact slide with table
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s); _accent(s)
    _box(s, Inches(0.7), Inches(0.5), Inches(11), Inches(0.4), 14, "QUANTIFIED IMPACT", GREEN, True)
    _box(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0), 34, "Measurable to the minute", TEXT, True)
    rows = [
        ("Metric", "Today", "With TORQ"),
        ("Time to diagnosis", "~30 min", "< 1 min (~95% faster)"),
        ("Dispatch lag", "15-60 min", "Instant, skill-matched"),
        ("Mean time to repair", "Baseline", "20-30% lower"),
        ("Knowledge retention", "Lost at turnover", "100% of logged fixes"),
        ("Supervisor admin", "Per work order", "~5 h saved / week"),
    ]
    tbl = s.shapes.add_table(len(rows), 3, Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.2)).table
    for c in range(3):
        tbl.columns[c].width = Inches(3.9)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN if r == 0 else CARD
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = val
            run.font.size = Pt(16); run.font.bold = (r == 0)
            run.font.color.rgb = BG if r == 0 else TEXT
            run.font.name = "Segoe UI"
    _box(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6), 14,
         "Worked example: ~$78,000/year of downtime avoided in a 10-fault/week plant, from a software-only tool.", MUTE)

    content(prs, "See it run", [
        "A machine faults live (MQTT). No human triggers it.",
        "Seconds later the diagnosis and work order appear in the approval queue.",
        "Supervisor taps Approve.",
        "A technician's phone buzzes with the full repair procedure, in Arabic.",
        "The dashboard shows time-to-diagnosis collapse from 30 min to under a minute.",
    ], kicker="Demo")

    content(prs, "Built to sell, built for $0", [
        "Buyer: maintenance managers of mid-size legacy plants (50-500 staff).",
        "Textile, agro-food, plastics, mechanical: Tunisia, North Africa, then global.",
        "Needs only the fault code and the paper-era manual, digitized once.",
        "Entire stack is free-tier: $0 infrastructure cost.",
        "Deploys in days, priced for SMEs, complements existing CMMS.",
    ], kicker="Business")

    content(prs, "Why we win", [
        "Enterprise CMMS (Maximo, SAP): they track work, they don't diagnose it.",
        "Modern CMMS (Fiix, UpKeep, MaintainX): still human-authored work orders.",
        "Generic AI chatbots: need manuals in the cloud, know nothing of the plant, act only when asked.",
        "TORQ is zero-touch, on-premise via MCP, and dispatches automatically.",
    ], kicker="Differentiation")

    # Close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.2), 40,
         "TORQ turns the most expensive 30 minutes\nin manufacturing into 30 seconds.", TEXT, True)
    _box(s, Inches(0.95), Inches(5.0), Inches(11), Inches(0.6), 20, "From fault code to fixed.", GREEN, True)

    prs.save(str(OUT))
    print(f"Deck written: {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
